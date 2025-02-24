import json
import pandas as pd
import pdfplumber
import re
from pptx import Presentation

class DataIngestion:
    """
    Handles reading data from various file formats:
      - JSON (datasets/dataset1.json)
      - CSV (datasets/dataset2.csv)
      - PDF (datasets/dataset3.pdf)
      - PPTX (datasets/dataset4.pptx)
    """
    def __init__(self, json_path='datasets/dataset1.json', csv_path='datasets/dataset2.csv',
                 pdf_path='datasets/dataset3.pdf', pptx_path='datasets/dataset4.pptx'):
        self.json_path = json_path
        self.csv_path = csv_path
        self.pdf_path = pdf_path
        self.pptx_path = pptx_path

    # JSON
    def load_json(self) -> dict:
        with open(self.json_path, 'r') as f:
            data = json.load(f)
        return data

    # CSV
    def load_csv(self) -> pd.DataFrame:
        df = pd.read_csv(self.csv_path)
        # Normalize column names (lowercase)
        df.columns = [col.strip().lower() for col in df.columns]
        return df

    # PDF
    def load_pdf(self) -> list:
        """
        Extract table data from the PDF.
        Returns a list of dictionaries (one per row).
        """
        all_rows = []
        with pdfplumber.open(self.pdf_path) as pdf:
            for page in pdf.pages:
                table = page.extract_table()
                if table:
                    headers = table[0]
                    for row in table[1:]:
                        if len(row) == len(headers):
                            row_dict = dict(zip(headers, row))
                            all_rows.append(row_dict)
        return all_rows

    # PPTX
    def load_pptx(self) -> dict:
        """
        Parse the PPTX slides individually. Returns a dictionary with:
          {
            "summary_metrics": {...},
            "quarterly_metrics": [...],
            "revenue_breakdown": {...}
          }
        """
        return self._parse_pptx_slides(self.pptx_path)

    def _load_pptx_all_slides(self, filepath):
        """
        Helper: returns a list of dicts, each with:
          - 'slide_index': the slide number
          - 'tables': list of 2D arrays (rows in each table)
          - 'texts': list of text from text-based shapes
        """
        prs = Presentation(filepath)
        slides_data = []

        for slide_index, slide in enumerate(prs.slides):
            slide_info = {
                "slide_index": slide_index,
                "tables": [],
                "texts": []
            }
            for shape in slide.shapes:
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        cells = [cell.text for cell in row.cells]
                        table_data.append(cells)
                    slide_info["tables"].append(table_data)
                elif shape.has_text_frame:
                    slide_info["texts"].append(shape.text)
            slides_data.append(slide_info)

        return slides_data

    def _parse_pptx_slides(self, filepath):
        """
        Parse each slide's content (text, tables) and apply different regex
        or logic depending on the slide index.
        """
        slides_data = self._load_pptx_all_slides(filepath)
        parsed_data = {
            "summary_metrics": {},
            "quarterly_metrics": [],
            "revenue_breakdown": {}
        }

        for slide_info in slides_data:
            slide_index = slide_info["slide_index"]
            text_content = "\n".join(slide_info["texts"])

            if slide_index == 0:
                revenue_match = re.search(r"Total Revenue:\s*\$([\d,]+)", text_content)
                if revenue_match:
                    val = revenue_match.group(1).replace(",", "")
                    parsed_data["summary_metrics"]["Total Revenue"] = int(val)

                memberships_match = re.search(r"Total Memberships Sold:\s*([\d,]+)", text_content)
                if memberships_match:
                    val = memberships_match.group(1).replace(",", "")
                    parsed_data["summary_metrics"]["Total Memberships Sold"] = int(val)

                location_match = re.search(r"Top Location:\s*([\w\s]+)", text_content)
                if location_match:
                    parsed_data["summary_metrics"]["Top Location"] = location_match.group(1).strip()

            elif slide_index == 1:
                if slide_info["tables"]:
                    table = slide_info["tables"][0]
                    headers = table[0]
                    for row in table[1:]:
                        row_dict = dict(zip(headers, row))
                        parsed_data["quarterly_metrics"].append(row_dict)

            elif slide_index == 2:
                revenue_breakdown_text = ""
                for text in slide_info["texts"]:
                    revenue_breakdown_text += text + "\n"
                pairs = re.findall(r"([\w\s]+):\s*(\d+)%", revenue_breakdown_text, re.IGNORECASE)
                for activity, pct in pairs:
                    parsed_data["revenue_breakdown"][activity.strip()] = int(pct)

        return parsed_data
